"""Generate CineSentiment final presentation as a .pptx file.

Matches the application's "video-club paper" design language:
paper backgrounds, dark ink text, red/green sentiment accents,
display + mono typographic hierarchy.

Output: docs/final_report/parts/CineSentiment_Presentation.pptx
Run:    python slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# --- Design tokens (mirror tailwind.config.ts) ---
PAPER = RGBColor(0xEF, 0xE5, 0xCD)          # #efe5cd
PAPER_2 = RGBColor(0xF7, 0xEF, 0xD9)        # #f7efd9
PAPER_DARK = RGBColor(0xE5, 0xD9, 0xBE)     # #e5d9be
INK = RGBColor(0x1B, 0x16, 0x12)            # #1b1612
INK_SOFT = RGBColor(0x3D, 0x34, 0x2A)       # #3d342a
RED = RGBColor(0x9B, 0x26, 0x14)            # #9b2614
RED_DEEP = RGBColor(0x7A, 0x1C, 0x10)
AMBER = RGBColor(0xE8, 0xA2, 0x3A)
KRAFT = RGBColor(0x8A, 0x6A, 0x3A)
GREEN = RGBColor(0x3D, 0x6B, 0x3A)

# Fonts (must exist on most macOS / Windows installs)
FONT_DISPLAY = "Impact"            # Anton/display fallback
FONT_BODY = "Georgia"              # serif body (Playfair fallback)
FONT_MONO = "Courier New"          # monospace meta

# --- Slide setup: 16:9 ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]


# ---- Helpers ----
def add_bg(slide, color=PAPER):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_box(slide, left, top, width, height, fill=PAPER_2,
            line_color=INK, line_width=1.0):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line_color
    box.line.width = Pt(line_width)
    return box


def add_text(slide, text, left, top, width, height,
             font=FONT_BODY, size=20, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, letter_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align

    if isinstance(text, list):
        lines = text
    else:
        lines = text.split("\n")

    for i, line in enumerate(lines):
        if i == 0:
            run = p.add_run()
        else:
            new_p = tf.add_paragraph()
            new_p.alignment = align
            run = new_p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        if letter_spacing is not None:
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(letter_spacing))
    return tb


def add_meta_label(slide, text, left, top, width=Inches(8), size=11):
    """Monospace caps label, à la the app's section headers."""
    return add_text(slide, text.upper(), left, top, width, Inches(0.35),
                    font=FONT_MONO, size=size, color=INK,
                    bold=True, letter_spacing=200)


def add_footer(slide, page_num, total):
    add_text(slide, "CINESENTIMENT — CMP 4418",
             Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             font=FONT_MONO, size=9, color=INK_SOFT, letter_spacing=200)
    add_text(slide, f"{page_num} / {total}",
             Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             font=FONT_MONO, size=9, color=INK_SOFT,
             align=PP_ALIGN.RIGHT, letter_spacing=200)


def add_divider(slide, left, top, width, color=INK):
    """Thick horizontal rule, brutalist style."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left, top, width, Pt(2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def add_perforation(slide, top=Inches(0.0)):
    """Decorative top stripe like a film reel header."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  0, top, W, Inches(0.18))
    rect.fill.solid()
    rect.fill.fore_color.rgb = INK
    rect.line.fill.background()


TOTAL_SLIDES = 13


# ===========================================================
# SLIDE 1 — Title
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)
add_perforation(s, top=Inches(0))
add_perforation(s, top=Inches(7.32))

add_meta_label(s, "REEL 01 · 35 MM · ACADEMIC PROJECT",
               Inches(0.5), Inches(0.55))

add_text(s, "CINESENTIMENT",
         Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.6),
         font=FONT_DISPLAY, size=110, color=INK, letter_spacing=-50)

add_text(s, "Fine-Tuned DistilBERT for Movie Review Sentiment Analysis",
         Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.6),
         font=FONT_BODY, size=22, color=INK_SOFT, italic=True)

add_divider(s, Inches(0.5), Inches(4.1), Inches(12.3))

add_text(s, "Esma Oruç      —      2022510140",
         Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.4),
         font=FONT_MONO, size=16, color=INK, letter_spacing=150)
add_text(s, "Ayan Abasova   —      2022510014",
         Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.4),
         font=FONT_MONO, size=16, color=INK, letter_spacing=150)

add_text(s, "CMP 4418  ·  FINAL PROJECT  ·  MAY 2026",
         Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
         font=FONT_MONO, size=12, color=INK_SOFT, letter_spacing=300)


# ===========================================================
# SLIDE 2 — The Problem
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)

add_meta_label(s, "REEL 02 · THE PROBLEM",
               Inches(0.5), Inches(0.4))
add_divider(s, Inches(0.5), Inches(0.9), Inches(12.3))

add_text(s, "Turning messy reviews\ninto structured signal.",
         Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.2),
         font=FONT_DISPLAY, size=64, color=INK)

add_text(s,
         "Movie reviews carry strong opinions but in a free, noisy form.",
         Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
         font=FONT_BODY, size=20, color=INK_SOFT, italic=True)

# Three feature boxes
y = Inches(4.55)
for i, (cap, body, accent) in enumerate([
    ("CHALLENGE", "Sarcasm, slang, comparatives, named entities.", RED),
    ("OPPORTUNITY", "IMDB 50k offers ground-truth labels at scale.", GREEN),
    ("APPROACH",  "Transformer encoder → end-to-end web product.", KRAFT),
]):
    x = Inches(0.5 + i * 4.18)
    add_box(s, x, y, Inches(4.0), Inches(2.0), fill=PAPER_2)
    add_box(s, x, y, Inches(0.2), Inches(2.0), fill=accent, line_color=accent)
    add_meta_label(s, cap, x + Inches(0.5), y + Inches(0.25), size=12)
    add_text(s, body, x + Inches(0.5), y + Inches(0.85),
             Inches(3.3), Inches(1.0),
             font=FONT_BODY, size=16, color=INK)

add_footer(s, 2, TOTAL_SLIDES)


# ===========================================================
# SLIDE 3 — The Pivot
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)
add_meta_label(s, "REEL 03 · THE PIVOT", Inches(0.5), Inches(0.4))
add_divider(s, Inches(0.5), Inches(0.9), Inches(12.3))

add_text(s, "From Twitter, to IMDB.",
         Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.1),
         font=FONT_DISPLAY, size=56, color=INK)

# Left column — original plan
add_box(s, Inches(0.5), Inches(2.7), Inches(6.0), Inches(4.0), fill=PAPER_2)
add_box(s, Inches(0.5), Inches(2.7), Inches(0.2), Inches(4.0), fill=RED, line_color=RED)
add_meta_label(s, "PART 1 · ORIGINAL PLAN",
               Inches(0.85), Inches(2.95), size=12)
add_text(s, "Live Twitter (X) sentiment on three movies side by side.",
         Inches(0.85), Inches(3.45), Inches(5.4), Inches(0.9),
         font=FONT_BODY, size=18, color=INK, italic=True)
add_text(s, [
    "✗  Free API retired in 2023.",
    "✗  Basic tier costs USD 100 / month.",
    "✗  Tweets have no ground-truth labels.",
    "✗  No reliable evaluation possible.",
], Inches(0.85), Inches(4.7), Inches(5.4), Inches(2.0),
   font=FONT_MONO, size=14, color=RED_DEEP)

# Right column — new plan
add_box(s, Inches(6.85), Inches(2.7), Inches(6.0), Inches(4.0), fill=PAPER_2)
add_box(s, Inches(6.85), Inches(2.7), Inches(0.2), Inches(4.0), fill=GREEN, line_color=GREEN)
add_meta_label(s, "FINAL · IMDB 50K",
               Inches(7.2), Inches(2.95), size=12)
add_text(s, "Labelled benchmark · Maas et al. 2011 · 25 k train + 25 k test.",
         Inches(7.2), Inches(3.45), Inches(5.4), Inches(0.9),
         font=FONT_BODY, size=18, color=INK, italic=True)
add_text(s, [
    "✓  Balanced 50/50 positive / negative.",
    "✓  Comparable with prior published work.",
    "✓  Confidence calibrated against truth.",
    "✓  Reproducible end-to-end.",
], Inches(7.2), Inches(4.7), Inches(5.4), Inches(2.0),
   font=FONT_MONO, size=14, color=GREEN)

add_footer(s, 3, TOTAL_SLIDES)


# ===========================================================
# SLIDE 4 — System Architecture
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)
add_meta_label(s, "REEL 04 · THE ARCHITECTURE",
               Inches(0.5), Inches(0.4))
add_divider(s, Inches(0.5), Inches(0.9), Inches(12.3))

add_text(s, "Four layers, replaceable.",
         Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0),
         font=FONT_DISPLAY, size=46, color=INK)

# Row of 4 boxes with arrows
labels = [
    ("TRAIN",   "Google Colab\nT4 GPU · 14 min", AMBER),
    ("STORE",   "Hugging Face Hub\nethsmaa/...", GREEN),
    ("SERVE",   "FastAPI Sidecar\nApple MPS · 50ms", RED),
    ("CONSUME", "Hono + tRPC\n→ React UI", KRAFT),
]
box_w = Inches(2.7)
y = Inches(3.0)
for i, (cap, body, accent) in enumerate(labels):
    x = Inches(0.5 + i * 3.15)
    add_box(s, x, y, box_w, Inches(1.9), fill=PAPER_2)
    add_box(s, x, y, Inches(0.18), Inches(1.9), fill=accent, line_color=accent)
    add_meta_label(s, cap, x + Inches(0.4), y + Inches(0.2), size=12)
    add_text(s, body, x + Inches(0.4), y + Inches(0.75),
             box_w - Inches(0.5), Inches(1.1),
             font=FONT_BODY, size=15, color=INK)
    if i < 3:
        # arrow
        arrow_x = x + box_w + Inches(0.05)
        arrow_box = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            arrow_x, y + Inches(0.75),
            Inches(0.35), Inches(0.4))
        arrow_box.fill.solid()
        arrow_box.fill.fore_color.rgb = INK
        arrow_box.line.fill.background()

# Bottom commentary
add_box(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4), fill=PAPER_DARK)
add_meta_label(s, "WHY THIS SHAPE",
               Inches(0.85), Inches(5.65), size=11)
add_text(s,
         "Each layer can be swapped without breaking the others: retrain the model, "
         "move serving to the cloud, replace the database — the contract is HTTP.",
         Inches(0.85), Inches(6.0), Inches(11.6), Inches(0.9),
         font=FONT_BODY, size=15, color=INK, italic=True)

add_footer(s, 4, TOTAL_SLIDES)


# ===========================================================
# SLIDE 5 — The Dataset
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)
add_meta_label(s, "REEL 05 · THE DATA",
               Inches(0.5), Inches(0.4))
add_divider(s, Inches(0.5), Inches(0.9), Inches(12.3))

add_text(s, "IMDB 50,000 reviews.",
         Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0),
         font=FONT_DISPLAY, size=46, color=INK)

# Stat blocks
stats = [
    ("25,000", "TRAIN", PAPER_2),
    ("25,000", "TEST",  PAPER_2),
    ("50 / 50", "BALANCE",  PAPER_2),
    ("256", "MAX TOKENS",   PAPER_2),
]
for i, (val, label, bg) in enumerate(stats):
    x = Inches(0.5 + i * 3.15)
    add_box(s, x, Inches(2.8), Inches(2.85), Inches(1.9), fill=bg)
    add_text(s, val, x, Inches(3.0), Inches(2.85), Inches(1.0),
             font=FONT_DISPLAY, size=48, color=INK,
             align=PP_ALIGN.CENTER)
    add_meta_label(s, label, x, Inches(4.0),
                   width=Inches(2.85), size=11)
    # center the label
    tf = s.shapes[-1].text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Sample reviews row
add_box(s, Inches(0.5), Inches(5.0), Inches(6.0), Inches(1.9), fill=PAPER_2)
add_box(s, Inches(0.5), Inches(5.0), Inches(0.2), Inches(1.9), fill=GREEN, line_color=GREEN)
add_meta_label(s, "POSITIVE",
               Inches(0.85), Inches(5.15), size=11)
add_text(s,
         "“A stunning, character-driven story with sharp dialogue and "
         "real emotional weight.”",
         Inches(0.85), Inches(5.55), Inches(5.4), Inches(1.3),
         font=FONT_BODY, size=14, color=INK, italic=True)

add_box(s, Inches(6.85), Inches(5.0), Inches(6.0), Inches(1.9), fill=PAPER_2)
add_box(s, Inches(6.85), Inches(5.0), Inches(0.2), Inches(1.9), fill=RED, line_color=RED)
add_meta_label(s, "NEGATIVE",
               Inches(7.2), Inches(5.15), size=11)
add_text(s,
         "“Two hours of overlong, predictable melodrama. Skip it.”",
         Inches(7.2), Inches(5.55), Inches(5.4), Inches(1.3),
         font=FONT_BODY, size=14, color=INK, italic=True)

add_footer(s, 5, TOTAL_SLIDES)


# ===========================================================
# SLIDE 6 — The Model
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)
add_meta_label(s, "REEL 06 · THE MODEL",
               Inches(0.5), Inches(0.4))
add_divider(s, Inches(0.5), Inches(0.9), Inches(12.3))

add_text(s, "DistilBERT.",
         Inches(0.5), Inches(1.2), Inches(7.0), Inches(1.0),
         font=FONT_DISPLAY, size=64, color=INK)
add_text(s, "smaller. faster. lighter.",
         Inches(0.5), Inches(2.4), Inches(7.0), Inches(0.5),
         font=FONT_BODY, size=20, color=INK_SOFT, italic=True)

# Right column — specs table
specs = [
    ("BASE MODEL",        "distilbert-base-uncased"),
    ("PARAMETERS",        "66 million"),
    ("LAYERS",            "6 encoder · 12 attention heads"),
    ("PRETRAIN",          "Wikipedia + BookCorpus"),
    ("CLASSIFIER HEAD",   "Linear, 2 labels on [CLS]"),
]
x0 = Inches(7.8)
y0 = Inches(1.3)
for i, (k, v) in enumerate(specs):
    y = y0 + Inches(i * 0.65)
    add_meta_label(s, k, x0, y, size=10)
    add_text(s, v, x0, y + Inches(0.25), Inches(5.0), Inches(0.4),
             font=FONT_MONO, size=13, color=INK)

# Bottom — sanh et al. quote
add_box(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.7), fill=PAPER_DARK)
add_meta_label(s, "LITERATURE",
               Inches(0.85), Inches(5.35), size=11)
add_text(s,
         "“DistilBERT is 40 % smaller, 60 % faster, and keeps "
         "≈ 97 % of BERT's downstream performance.”",
         Inches(0.85), Inches(5.75), Inches(11.6), Inches(0.6),
         font=FONT_BODY, size=18, color=INK, italic=True)
add_text(s, "— Sanh, Debut, Chaumond, Wolf · 2019",
         Inches(0.85), Inches(6.35), Inches(11.6), Inches(0.5),
         font=FONT_MONO, size=11, color=INK_SOFT)

add_footer(s, 6, TOTAL_SLIDES)


# ===========================================================
# SLIDE 7 — Training Configuration
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)
add_meta_label(s, "REEL 07 · THE PRODUCTION",
               Inches(0.5), Inches(0.4))
add_divider(s, Inches(0.5), Inches(0.9), Inches(12.3))

add_text(s, "14 minutes on a T4.",
         Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0),
         font=FONT_DISPLAY, size=52, color=INK)

# Two-column hyperparameter list
left_specs = [
    ("EPOCHS",            "3"),
    ("BATCH SIZE",        "16 per device"),
    ("LEARNING RATE",     "2 × 10⁻⁵"),
    ("OPTIMIZER",         "AdamW"),
]
right_specs = [
    ("MAX SEQ LENGTH",    "256 tokens"),
    ("PRECISION",         "FP16 mixed"),
    ("HARDWARE",          "NVIDIA T4 · Google Colab"),
    ("WALL-CLOCK",        "858.6 s · ≈14.3 min"),
]
for col, specs in enumerate([left_specs, right_specs]):
    x0 = Inches(0.5 + col * 6.4)
    y0 = Inches(2.7)
    for i, (k, v) in enumerate(specs):
        y = y0 + Inches(i * 0.85)
        add_box(s, x0, y, Inches(6.0), Inches(0.7), fill=PAPER_2)
        add_meta_label(s, k, x0 + Inches(0.25), y + Inches(0.12), size=10)
        add_text(s, v, x0 + Inches(2.4), y + Inches(0.15),
                 Inches(3.4), Inches(0.5),
                 font=FONT_MONO, size=14, color=INK, bold=True)

add_footer(s, 7, TOTAL_SLIDES)


# ===========================================================
# SLIDE 8 — Three-Class Threshold
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)
add_meta_label(s, "REEL 08 · DIRECTOR'S COMMENTARY",
               Inches(0.5), Inches(0.4))
add_divider(s, Inches(0.5), Inches(0.9), Inches(12.3))

add_text(s, "Binary model, three buckets.",
         Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0),
         font=FONT_DISPLAY, size=44, color=INK)

add_text(s,
         "The classifier outputs p⁺ ∈ [0, 1]. The UI shows three classes.",
         Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.5),
         font=FONT_BODY, size=18, color=INK_SOFT, italic=True)

# Threshold visualization
bar_x = Inches(0.5)
bar_y = Inches(3.5)
bar_h = Inches(0.7)
bar_total = Inches(12.3)

# Three coloured segments: neg [0, 0.35], neu [0.35, 0.65], pos [0.65, 1.0]
neg_w = Emu(int(bar_total * 0.35))
neu_w = Emu(int(bar_total * 0.30))
pos_w = Emu(int(bar_total * 0.35))

add_box(s, bar_x, bar_y, neg_w, bar_h, fill=RED, line_color=INK)
add_box(s, bar_x + neg_w, bar_y, neu_w, bar_h, fill=KRAFT, line_color=INK)
add_box(s, bar_x + neg_w + neu_w, bar_y, pos_w, bar_h, fill=GREEN, line_color=INK)

# Segment labels (inside bars, white)
add_text(s, "NEGATIVE", bar_x, bar_y, neg_w, bar_h,
         font=FONT_MONO, size=14, color=PAPER, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         letter_spacing=200)
add_text(s, "NEUTRAL", bar_x + neg_w, bar_y, neu_w, bar_h,
         font=FONT_MONO, size=14, color=PAPER, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         letter_spacing=200)
add_text(s, "POSITIVE", bar_x + neg_w + neu_w, bar_y, pos_w, bar_h,
         font=FONT_MONO, size=14, color=PAPER, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         letter_spacing=200)

# Threshold markers
add_text(s, "p⁺ = 0.00", bar_x, bar_y + bar_h + Inches(0.1),
         Inches(1.5), Inches(0.4),
         font=FONT_MONO, size=11, color=INK_SOFT)
add_text(s, "p⁺ = 0.35", bar_x + neg_w - Inches(0.75),
         bar_y + bar_h + Inches(0.1),
         Inches(1.5), Inches(0.4),
         font=FONT_MONO, size=11, color=INK, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, "p⁺ = 0.65",
         bar_x + neg_w + neu_w - Inches(0.75),
         bar_y + bar_h + Inches(0.1),
         Inches(1.5), Inches(0.4),
         font=FONT_MONO, size=11, color=INK, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, "p⁺ = 1.00", bar_x + bar_total - Inches(1.5),
         bar_y + bar_h + Inches(0.1),
         Inches(1.5), Inches(0.4),
         font=FONT_MONO, size=11, color=INK_SOFT,
         align=PP_ALIGN.RIGHT)

# Honesty note
add_box(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4), fill=PAPER_DARK)
add_meta_label(s, "CALLED OUT IN THE MODEL CARD",
               Inches(0.85), Inches(5.65), size=10)
add_text(s,
         "Heuristic, not learned. Documented openly in both the source "
         "code and the in-app model page.",
         Inches(0.85), Inches(6.0), Inches(11.6), Inches(0.9),
         font=FONT_BODY, size=15, color=INK, italic=True)

add_footer(s, 8, TOTAL_SLIDES)


# ===========================================================
# SLIDE 9 — Strict Mode
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)
add_meta_label(s, "REEL 09 · STRICT MODE",
               Inches(0.5), Inches(0.4))
add_divider(s, Inches(0.5), Inches(0.9), Inches(12.3))

add_text(s, "No silent fallbacks.",
         Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0),
         font=FONT_DISPLAY, size=58, color=INK)

add_text(s,
         "Every prediction shown in the UI must come from the trained model.",
         Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.5),
         font=FONT_BODY, size=20, color=INK_SOFT, italic=True)

# Two boxes — before and after
add_box(s, Inches(0.5), Inches(3.2), Inches(6.0), Inches(3.5), fill=PAPER_2)
add_box(s, Inches(0.5), Inches(3.2), Inches(0.2), Inches(3.5), fill=RED, line_color=RED)
add_meta_label(s, "BEFORE", Inches(0.85), Inches(3.4), size=11)
add_text(s, "Quiet fallback to a lexicon stand-in when FastAPI is down. "
            "Several hours of incorrect labels could go unnoticed.",
         Inches(0.85), Inches(3.85), Inches(5.4), Inches(2.8),
         font=FONT_BODY, size=16, color=INK)

add_box(s, Inches(6.85), Inches(3.2), Inches(6.0), Inches(3.5), fill=PAPER_2)
add_box(s, Inches(6.85), Inches(3.2), Inches(0.2), Inches(3.5), fill=GREEN, line_color=GREEN)
add_meta_label(s, "AFTER", Inches(7.2), Inches(3.4), size=11)
add_text(s, "BERT_STRICT=true. Sidecar outage → loud error. "
            "Database carries one model version: distilbert-imdb-v1.",
         Inches(7.2), Inches(3.85), Inches(5.4), Inches(2.8),
         font=FONT_BODY, size=16, color=INK)

add_footer(s, 9, TOTAL_SLIDES)


# ===========================================================
# SLIDE 10 — Results headline
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)
add_meta_label(s, "REEL 10 · THE NUMBERS",
               Inches(0.5), Inches(0.4))
add_divider(s, Inches(0.5), Inches(0.9), Inches(12.3))

# Hero metric
add_text(s, "91.51 %",
         Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.2),
         font=FONT_DISPLAY, size=240, color=INK,
         align=PP_ALIGN.CENTER, letter_spacing=-80)

add_text(s, "ACCURACY  ON THE  IMDB  TEST SET",
         Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
         font=FONT_MONO, size=18, color=INK, letter_spacing=400,
         align=PP_ALIGN.CENTER)

# Secondary metrics row
for i, (val, label) in enumerate([
    ("0.9151", "WEIGHTED F1"),
    ("0.9151", "PRECISION"),
    ("0.9151", "RECALL"),
    ("25,000", "TEST SET"),
]):
    x = Inches(0.5 + i * 3.15)
    add_box(s, x, Inches(5.6), Inches(2.85), Inches(1.2), fill=PAPER_2)
    add_text(s, val, x, Inches(5.7), Inches(2.85), Inches(0.6),
             font=FONT_DISPLAY, size=28, color=INK,
             align=PP_ALIGN.CENTER)
    add_meta_label(s, label, x, Inches(6.3),
                   width=Inches(2.85), size=10)
    tf = s.shapes[-1].text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text(s,
         "Target was 85 %. We beat it by more than six points.",
         Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
         font=FONT_BODY, size=14, color=INK_SOFT, italic=True,
         align=PP_ALIGN.CENTER)


# ===========================================================
# SLIDE 11 — Confusion Matrix + per-class
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)
add_meta_label(s, "REEL 11 · ERROR STRUCTURE",
               Inches(0.5), Inches(0.4))
add_divider(s, Inches(0.5), Inches(0.9), Inches(12.3))

add_text(s, "Symmetric errors. Balanced classifier.",
         Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0),
         font=FONT_DISPLAY, size=38, color=INK)

# Confusion matrix table (left)
matrix_x = Inches(0.5)
matrix_y = Inches(2.6)
cell_w = Inches(2.2)
cell_h = Inches(1.0)

# Header row
add_text(s, "PREDICTED →",
         matrix_x, matrix_y - Inches(0.4),
         Inches(6.6), Inches(0.4),
         font=FONT_MONO, size=11, color=INK_SOFT,
         align=PP_ALIGN.CENTER, letter_spacing=200)

# Empty top-left
add_text(s, "", matrix_x, matrix_y, cell_w, cell_h)
add_box(s, matrix_x + cell_w, matrix_y, cell_w, cell_h,
        fill=PAPER_DARK)
add_text(s, "NEGATIVE", matrix_x + cell_w, matrix_y,
         cell_w, cell_h, font=FONT_MONO, size=11,
         color=INK, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         letter_spacing=200)
add_box(s, matrix_x + cell_w*2, matrix_y, cell_w, cell_h,
        fill=PAPER_DARK)
add_text(s, "POSITIVE", matrix_x + cell_w*2, matrix_y,
         cell_w, cell_h, font=FONT_MONO, size=11,
         color=INK, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         letter_spacing=200)

# Row labels and values
rows = [
    ("ACTUAL NEG", "11,394", "1,106", GREEN, RED),
    ("ACTUAL POS", "1,017",  "11,483", RED, GREEN),
]
for i, (lbl, v1, v2, c1, c2) in enumerate(rows):
    y = matrix_y + cell_h * (i + 1)
    add_box(s, matrix_x, y, cell_w, cell_h, fill=PAPER_DARK)
    add_text(s, lbl, matrix_x, y, cell_w, cell_h,
             font=FONT_MONO, size=11, color=INK, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             letter_spacing=200)
    for j, (v, c) in enumerate([(v1, c1), (v2, c2)]):
        cx = matrix_x + cell_w * (j + 1)
        add_box(s, cx, y, cell_w, cell_h, fill=PAPER_2)
        add_text(s, v, cx, y, cell_w, cell_h,
                 font=FONT_DISPLAY, size=26, color=c,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Per-class table on the right
x0 = Inches(7.5)
y0 = Inches(2.6)
add_meta_label(s, "PER-CLASS METRICS", x0, y0 - Inches(0.4),
               size=11)
classes = [
    ("NEGATIVE", "0.918", "0.912", "0.915", "12,500", RED),
    ("POSITIVE", "0.912", "0.919", "0.915", "12,500", GREEN),
]
header_labels = ["CLASS", "PREC", "RECALL", "F1", "SUPPORT"]
col_w = Inches(1.05)

# Header
for j, h in enumerate(header_labels):
    cx = x0 + col_w * j
    add_box(s, cx, y0, col_w, Inches(0.6), fill=PAPER_DARK)
    add_text(s, h, cx, y0, col_w, Inches(0.6),
             font=FONT_MONO, size=10, color=INK, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             letter_spacing=200)

# Rows
for i, row in enumerate(classes):
    name, p, r, f1, sup, c = row
    y = y0 + Inches(0.6) + Inches(i * 0.7)
    # Accent stripe
    add_box(s, x0, y, Inches(0.15), Inches(0.7), fill=c, line_color=c)
    # Cells
    for j, v in enumerate([name, p, r, f1, sup]):
        cx = x0 + col_w * j
        add_box(s, cx, y, col_w, Inches(0.7), fill=PAPER_2)
        add_text(s, v, cx, y, col_w, Inches(0.7),
                 font=FONT_MONO, size=12, color=INK, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s,
         "Both classes within one F1 percentage point of each other. "
         "No re-weighting needed.",
         Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
         font=FONT_BODY, size=15, color=INK_SOFT, italic=True,
         align=PP_ALIGN.CENTER)

add_footer(s, 11, TOTAL_SLIDES)


# ===========================================================
# SLIDE 12 — Live System / Demo
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)
add_meta_label(s, "REEL 12 · OPEN PRODUCTION",
               Inches(0.5), Inches(0.4))
add_divider(s, Inches(0.5), Inches(0.9), Inches(12.3))

add_text(s, "Live demo.",
         Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0),
         font=FONT_DISPLAY, size=72, color=INK)

# Three columns of links
links = [
    ("MODEL",
     "huggingface.co/\nethsmaa/\ncinesentiment-\ndistilbert-imdb",
     GREEN),
    ("CODE",
     "github.com/\nethsmaa/\nmovie-sentiment",
     KRAFT),
    ("APP",
     "localhost:5173\n\n· grid\n· analyze\n· model page",
     RED),
]
for i, (cap, body, accent) in enumerate(links):
    x = Inches(0.5 + i * 4.18)
    add_box(s, x, Inches(2.8), Inches(4.0), Inches(3.7), fill=PAPER_2)
    add_box(s, x, Inches(2.8), Inches(0.2), Inches(3.7),
            fill=accent, line_color=accent)
    add_meta_label(s, cap, x + Inches(0.5), Inches(3.0), size=12)
    add_text(s, body, x + Inches(0.5), Inches(3.7),
             Inches(3.3), Inches(2.7),
             font=FONT_MONO, size=14, color=INK)

add_text(s,
         "$ pnpm dev   →   api · web · ml + Postgres up in 10 seconds",
         Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
         font=FONT_MONO, size=14, color=INK_SOFT,
         align=PP_ALIGN.CENTER, letter_spacing=100)


# ===========================================================
# SLIDE 13 — Thank you / Q&A
# ===========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, PAPER)
add_perforation(s, top=Inches(0))
add_perforation(s, top=Inches(7.32))

add_meta_label(s, "REEL 13 · END CREDITS",
               Inches(0.5), Inches(0.55))

add_text(s, "Thank you.",
         Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.0),
         font=FONT_DISPLAY, size=180, color=INK,
         align=PP_ALIGN.CENTER, letter_spacing=-80)

add_text(s, "questions  /  discussion",
         Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
         font=FONT_BODY, size=24, color=INK_SOFT,
         align=PP_ALIGN.CENTER, italic=True)

add_divider(s, Inches(4.5), Inches(5.5), Inches(4.3))

add_text(s, "Esma Oruç · 2022510140",
         Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.4),
         font=FONT_MONO, size=14, color=INK,
         align=PP_ALIGN.CENTER, letter_spacing=200)
add_text(s, "Ayan Abasova · 2022510014",
         Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
         font=FONT_MONO, size=14, color=INK,
         align=PP_ALIGN.CENTER, letter_spacing=200)

add_text(s, "CINESENTIMENT  ·  CMP 4418  ·  MAY 2026",
         Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
         font=FONT_MONO, size=10, color=INK_SOFT,
         align=PP_ALIGN.CENTER, letter_spacing=400)


# ===========================================================
# Save
# ===========================================================
out = "CineSentiment_Presentation.pptx"
prs.save(out)
print(f"Saved {out}  ·  {TOTAL_SLIDES} slides")
